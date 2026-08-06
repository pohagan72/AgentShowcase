# features/pii_redaction/routes.py
from flask import (
    Blueprint, render_template, request, flash, current_app, url_for, g, redirect, send_file, session
)
import os
import io
import uuid
from werkzeug.utils import secure_filename
from docx import Document
from pptx import Presentation
import logging

# Shared rate limiter
from extensions import limiter

# Define the Blueprint
bp = Blueprint('pii_redaction', __name__)

# Entity sets per redaction mode. Presidio's default output (no entities= filter)
# is too aggressive on CV-shaped documents — it flags cities, employment dates,
# nationalities, etc. We opt in per mode instead.
#
# PII_ONLY: things that identify a specific person or account (names, contact,
# IDs, cards, URLs). Deliberately excludes LOCATION and DATE_TIME so cities and
# "five years experience" survive.
# FULL_ANON: adds LOCATION, DATE_TIME, NRP, and ORGANIZATION for GDPR-style
# anonymisation of CVs / memos where employers, cities, dates, and nationality
# should also be blanked. ORGANIZATION requires the custom analyzer built by
# build_analyzer() below — Presidio's default English registry strips ORG out
# of the SpacyRecognizer's supported_entities.
_PII_ONLY_ENTITIES = [
    "PERSON", "EMAIL_ADDRESS", "PHONE_NUMBER",
    "US_SSN", "CREDIT_CARD", "IBAN_CODE",
    "US_PASSPORT", "US_DRIVER_LICENSE",
    "URL",
]
_FULL_ANON_ENTITIES = _PII_ONLY_ENTITIES + [
    "LOCATION", "DATE_TIME", "NRP", "ORGANIZATION",
]
# 0.4 catches phone numbers (~0.4) and US passports (~0.45) in weaker context
# while still filtering the sub-0.4 noise on ordinary title-cased words like
# "Microsoft" or "Foundation Certificate". Presidio's zero-threshold default
# is what caused the CV over-redaction in the first place.
_MIN_SCORE = 0.4

# ■ (U+25A0 Black Square) instead of █ (U+2588 Full Block). U+2588 has poor
# font coverage — Word substitutes glyphs from a fallback font and the output
# renders as a mess of ù / empty rectangles in typical CV title fonts. U+25A0
# ships in every mainstream font and reads as the same "redacted bar" intent.
_REDACT_CHAR = "■"


def build_analyzer():
    """Construct Presidio's AnalyzerEngine with ORGANIZATION enabled.

    Presidio's default English SpacyRecognizer instance ships with
    supported_entities = [DATE_TIME, PERSON, LOCATION, NRP] — ORGANIZATION is
    stripped out even though spaCy's NER emits ORG labels for free. This
    factory swaps in a SpacyRecognizer configured to keep the ORG channel so
    "Full anonymisation" mode can actually blank employer names.

    Callers must provide the nlp_engine (the app builds one lazily at startup).
    Returns an AnalyzerEngine ready to hand to redact_*_document_pii.
    """
    from presidio_analyzer import AnalyzerEngine, RecognizerRegistry
    from presidio_analyzer.predefined_recognizers import SpacyRecognizer
    from presidio_analyzer.nlp_engine import NlpEngineProvider

    provider = NlpEngineProvider(nlp_configuration={
        "nlp_engine_name": "spacy",
        "models": [{"lang_code": "en", "model_name": "en_core_web_lg"}],
    })
    nlp_engine = provider.create_engine()

    registry = RecognizerRegistry()
    # Load the standard recognizer set (EmailRecognizer, PhoneRecognizer, etc.),
    # then replace the default SpacyRecognizer with one that keeps ORGANIZATION.
    registry.load_predefined_recognizers(languages=["en"], nlp_engine=nlp_engine)
    registry.remove_recognizer("SpacyRecognizer")
    registry.add_recognizer(SpacyRecognizer(
        supported_language="en",
        supported_entities=["DATE_TIME", "NRP", "LOCATION", "PERSON", "ORGANIZATION"],
    ))

    return AnalyzerEngine(
        registry=registry,
        nlp_engine=nlp_engine,
        supported_languages=["en"],
    )


def entities_for_mode(mode):
    """Map the client-supplied mode string to a Presidio entity allowlist.

    Unknown / missing mode falls back to PII_ONLY so a broken form field can't
    accidentally trigger the more aggressive anonymisation pass.
    """
    if mode == "full_anon":
        return _FULL_ANON_ENTITIES
    return _PII_ONLY_ENTITIES


def allowed_file_pii(filename):
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in current_app.config.get('PII_ALLOWED_EXTENSIONS', {'docx', 'pptx'})

def apply_redaction_to_text(text, analysis_results):
    """
    Helper function to apply character-level redaction to a string.
    Returns the redacted string.
    """
    if not analysis_results:
        return text
    
    # Convert to list for mutable character replacement
    text_chars = list(text)
    
    for res in analysis_results:
        start = res.start
        end = res.end
        # Ensure indices are within bounds
        start = max(0, start)
        end = min(len(text_chars), end)
        
        # Replace characters with block symbol
        for i in range(start, end):
            text_chars[i] = _REDACT_CHAR
            
    return "".join(text_chars)

def redact_runs_in_paragraph(paragraph, analyzer, entities=None):
    """
    Analyzes a paragraph and redacts PII within its runs strictly where overlap occurs.
    entities: iterable of Presidio entity names to detect; None = PII-only default.
    """
    text = paragraph.text
    if not text.strip():
        return False

    if entities is None:
        entities = _PII_ONLY_ENTITIES

    try:
        results = analyzer.analyze(
            text=text, language='en', entities=entities, score_threshold=_MIN_SCORE
        )
    except Exception as e:
        logging.error(f"Error analyzing paragraph text: {e}")
        return False

    if not results:
        return False

    current_offset = 0
    redaction_occurred = False

    # Iterate through runs and map global PII indices to local run indices
    for run in paragraph.runs:
        run_text = run.text
        run_len = len(run_text)
        
        # Skip empty runs
        if run_len == 0:
            continue
            
        run_start = current_offset
        run_end = current_offset + run_len
        
        # We process the run if it overlaps with any PII result
        # We rebuild the run string char by char to handle multiple PIIs in one run
        new_run_chars = list(run_text)
        run_modified = False

        for res in results:
            pii_start = res.start
            pii_end = res.end
            
            # Check overlap
            overlap_start = max(run_start, pii_start)
            overlap_end = min(run_end, pii_end)
            
            if overlap_start < overlap_end:
                # Calculate local indices relative to this run
                local_start = overlap_start - run_start
                local_end = overlap_end - run_start
                
                # Replace specific characters
                for i in range(local_start, local_end):
                    new_run_chars[i] = _REDACT_CHAR
                
                run_modified = True
                redaction_occurred = True

        if run_modified:
            run.text = "".join(new_run_chars)
        
        current_offset += run_len
    
    return redaction_occurred

def redact_word_document_pii(file_stream, analyzer, entities=None):
    try:
        document = Document(file_stream)
        redacted_count = 0
        logging.info(f"[{g.request_id if hasattr(g, 'request_id') else 'PII_REDACT'}] Starting Word document redaction.")

        # 1. Process Paragraphs
        for para in document.paragraphs:
            if redact_runs_in_paragraph(para, analyzer, entities=entities):
                redacted_count += 1

        # 2. Process Tables
        for table in document.tables:
            for row in table.rows:
                for cell in row.cells:
                    for para in cell.paragraphs:
                        if redact_runs_in_paragraph(para, analyzer, entities=entities):
                            redacted_count += 1
        
        logging.info(f"[{g.request_id if hasattr(g, 'request_id') else 'PII_REDACT'}] Modified approx {redacted_count} paragraphs/cells in Word document.")
        
        output_stream = io.BytesIO()
        document.save(output_stream)
        output_stream.seek(0)
        return output_stream
    except Exception as e:
        logging.error(f"[{g.request_id if hasattr(g, 'request_id') else 'PII_REDACT'}] Error processing Word document for PII: {e}", exc_info=True)
        return None

def redact_powerpoint_document_pii(file_stream, analyzer, entities=None):
    if entities is None:
        entities = _PII_ONLY_ENTITIES
    try:
        presentation = Presentation(file_stream)
        redacted_count = 0
        req_id_tag = g.request_id if hasattr(g, 'request_id') else 'PII_REDACT_PPTX'
        logging.info(f"[{req_id_tag}] Starting PowerPoint document redaction.")

        for i, slide in enumerate(presentation.slides):
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                
                text_frame = shape.text_frame
                for para in text_frame.paragraphs:
                    text = para.text
                    if not text.strip():
                        continue

                    try:
                        results = analyzer.analyze(
                            text=text, language='en',
                            entities=entities, score_threshold=_MIN_SCORE,
                        )
                    except Exception as e:
                        logging.error(f"[{req_id_tag}] Error analyzing PPTX paragraph: {e}")
                        continue

                    if not results:
                        continue

                    current_offset = 0
                    
                    # Similar logic to Word: Map global paragraph offsets to run offsets
                    for run in para.runs:
                        # Some runs in PPTX might be empty or None
                        if not hasattr(run, 'text') or not run.text:
                            continue

                        run_text = run.text
                        run_len = len(run_text)
                        
                        run_start = current_offset
                        run_end = current_offset + run_len
                        
                        new_run_chars = list(run_text)
                        run_modified = False

                        for res in results:
                            pii_start = res.start
                            pii_end = res.end
                            
                            overlap_start = max(run_start, pii_start)
                            overlap_end = min(run_end, pii_end)
                            
                            if overlap_start < overlap_end:
                                local_start = overlap_start - run_start
                                local_end = overlap_end - run_start
                                
                                for k in range(local_start, local_end):
                                    new_run_chars[k] = _REDACT_CHAR
                                
                                run_modified = True
                                redacted_count += 1
                        
                        if run_modified:
                            run.text = "".join(new_run_chars)
                        
                        current_offset += run_len

        logging.info(f"[{req_id_tag}] Redacted content in approx {redacted_count} runs in PowerPoint document.")
        
        output_stream = io.BytesIO()
        presentation.save(output_stream)
        output_stream.seek(0)
        return output_stream
    except Exception as e:
        logging.error(f"[{req_id_tag}] Error processing PowerPoint document for PII: {e}", exc_info=True)
        return None

@bp.route("/process/pii_redaction/redact", methods=["POST"])
@limiter.limit("10 per hour; 2 per minute")
def process_redact():
    g.request_id = uuid.uuid4().hex
    logging.info(f"Processing PII redaction request {g.request_id}...")

    context = {
        "redacted_file_url": None,
        "original_filename": None,
        "presidio_available": current_app.config.get('PRESIDIO_ANALYZER_AVAILABLE', False),
        "gcs_available": current_app.config.get('GCS_AVAILABLE', False),
        "hx_target_is_result": True
    }

    if not context["presidio_available"]:
        flash("PII Redaction service (Presidio Analyzer) is not available.", "error")
        logging.error(f"[{g.request_id}] Presidio Analyzer not available.")
        return render_template("pii_redaction/templates/pii_redaction_content.html", **context)

    if not context["gcs_available"]:
        flash("Cloud Storage is not available. Cannot store redacted file.", "error")
        logging.error(f"[{g.request_id}] GCS not available for PII redaction output.")
        return render_template("pii_redaction/templates/pii_redaction_content.html", **context)

    if 'file_to_redact' not in request.files:
        flash('No file part selected for redaction.', 'error')
        return render_template("pii_redaction/templates/pii_redaction_content.html", **context)

    file = request.files['file_to_redact']
    if file.filename == '':
        flash('No file selected for redaction.', 'error')
        return render_template("pii_redaction/templates/pii_redaction_content.html", **context)

    if file and allowed_file_pii(file.filename):
        original_filename = secure_filename(file.filename)
        context["original_filename"] = original_filename
        file_ext = original_filename.rsplit('.', 1)[1].lower()

        mode = request.form.get("mode", "full_anon")
        entities = entities_for_mode(mode)
        logging.info(f"[{g.request_id}] File received: {original_filename} (Type: {file_ext}, Mode: {mode})")

        file_stream = io.BytesIO(file.read())
        output_stream = None

        analyzer = current_app.presidio_analyzer
        gcs_bucket = current_app.gcs_bucket

        try:
            if file_ext == 'docx':
                output_stream = redact_word_document_pii(file_stream, analyzer, entities=entities)
            elif file_ext == 'pptx':
                output_stream = redact_powerpoint_document_pii(file_stream, analyzer, entities=entities)
            
            if output_stream:
                redacted_gcs_path = f"pii_redaction_results/{g.request_id}/redacted_{original_filename}"
                redacted_blob = gcs_bucket.blob(redacted_gcs_path)
                
                output_stream.seek(0)
                redacted_blob.upload_from_file(output_stream)
                logging.info(f"[{g.request_id}] Redacted file '{original_filename}' uploaded to GCS: gs://{gcs_bucket.name}/{redacted_gcs_path}")

                session_file_id = f"redacted_gcs_{g.request_id}"
                session[session_file_id] = {
                    'gcs_path': redacted_gcs_path,
                    'filename': f"redacted_{original_filename}",
                    'mimetype': f'application/vnd.openxmlformats-officedocument.{"wordprocessingml.document" if file_ext == "docx" else "presentationml.presentation"}'
                }
                context["redacted_file_url"] = url_for('pii_redaction.download_redacted_file_pii', file_id=session_file_id)
                flash(f"Document '{original_filename}' processed for PII redaction. Click link to download.", "success")
            else:
                flash(f"PII redaction process failed for '{original_filename}'. Output stream was empty. Check logs.", "error")
                logging.error(f"[{g.request_id}] Redaction output_stream was None for {original_filename}.")

        except Exception as e:
            flash(f'An error occurred during PII redaction of "{original_filename}": {str(e)}', "error")
            logging.error(f"[{g.request_id}] Error during PII redaction route for '{original_filename}': {e}", exc_info=True)
        finally:
            file_stream.close()
            if output_stream:
                output_stream.close()
    else:
        flash('Invalid file type for PII redaction. Please upload .docx or .pptx files.', 'error')
        logging.warning(f"[{g.request_id}] Invalid file type: {file.filename if file else 'No file'}")

    return render_template("pii_redaction/templates/pii_redaction_content.html", **context)

@bp.route("/process/pii_redaction/download/<file_id>")
def download_redacted_file_pii(file_id):
    gcs_bucket = current_app.gcs_bucket
    if not current_app.config.get('GCS_AVAILABLE') or not gcs_bucket:
        flash("Cloud Storage is not available for download.", "error")
        return redirect(url_for('main.index', feature_key='pii_redaction'))

    file_info = session.get(file_id)
    
    if not file_info or 'gcs_path' not in file_info:
        flash("Redacted file information not found or link expired.", "error")
        session.pop(file_id, None)
        return redirect(url_for('main.index', feature_key='pii_redaction'))

    gcs_path = file_info['gcs_path']
    filename_for_download = file_info.get('filename', 'redacted_file')
    mimetype = file_info.get('mimetype', 'application/octet-stream')

    try:
        logging.info(f"Attempting to download redacted file from GCS: gs://{gcs_bucket.name}/{gcs_path}")
        blob = gcs_bucket.blob(gcs_path)
        
        if not blob.exists():
            logging.error(f"Blob not found at GCS path: {gcs_path}")
            flash("Error: Redacted file no longer found in cloud storage (it may have expired).", "error")
            session.pop(file_id, None)
            return redirect(url_for('main.index', feature_key='pii_redaction'))

        output_stream = io.BytesIO()
        blob.download_to_file(output_stream)
        output_stream.seek(0)
        
        logging.info(f"Successfully downloaded '{filename_for_download}' from GCS for serving.")
        
        session.pop(file_id, None) 

        return send_file(
            output_stream,
            as_attachment=True,
            download_name=filename_for_download,
            mimetype=mimetype
        )
    except Exception:
        logging.error(f"GCSNotFound: Blob not found at GCS path: {gcs_path} during download attempt.")
        flash("Error: Redacted file not found in cloud storage during download attempt.", "error")
        session.pop(file_id, None)
        return redirect(url_for('main.index', feature_key='pii_redaction'))
    except Exception as e:
        logging.error(f"Error serving redacted file from GCS '{gcs_path}': {e}", exc_info=True)
        flash(f"An error occurred while trying to serve the redacted file: {str(e)}", "error")
        session.pop(file_id, None)
        return redirect(url_for('main.index', feature_key='pii_redaction'))