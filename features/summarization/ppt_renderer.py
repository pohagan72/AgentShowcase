# features/summarization/ppt_renderer.py
import logging
from io import BytesIO
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import MSO_ANCHOR, PP_PARAGRAPH_ALIGNMENT
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE

# --- Style Definitions ---
TEMPLATES = {
    'professional': { 
        'bg_color': RGBColor(255, 255, 255), 
        'title_color': RGBColor(0, 51, 102),     # Navy Blue
        'accent_color': RGBColor(0, 102, 204),   # Bright Blue
        'text_color': RGBColor(60, 60, 60), 
        'panel_bg': RGBColor(240, 245, 250),     # Light Blue Grey
        'panel_border': RGBColor(200, 200, 200)
    },
    'creative': { 
        'bg_color': RGBColor(250, 250, 250), 
        'title_color': RGBColor(128, 0, 0),      # Maroon
        'accent_color': RGBColor(204, 102, 0),   # Burnt Orange
        'text_color': RGBColor(50, 50, 50),
        'panel_bg': RGBColor(255, 245, 235),
        'panel_border': RGBColor(230, 200, 180)
    },
    'minimalist': {
        'bg_color': RGBColor(255, 255, 255),
        'title_color': RGBColor(0, 0, 0),        # Black
        'accent_color': RGBColor(100, 100, 100), # Grey
        'text_color': RGBColor(30, 30, 30),
        'panel_bg': RGBColor(245, 245, 245),
        'panel_border': RGBColor(220, 220, 220)
    }
}

DEFAULT_TEMPLATE_NAME = 'professional'

def add_insight_panel(slide, text, template, left, top, width, height):
    """
    Adds a styled 'Strategic Insight' sidebar to the right side of the slide.
    """
    if not text or text.lower() in ["n/a", "none", ""]: 
        return

    # 1. Background Box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = template['panel_bg']
    shape.line.color.rgb = template['panel_border']
    shape.line.width = Pt(1.5)

    # 2. Header Box ("Strategic Insight")
    header_height = Inches(0.6)
    header_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, header_height)
    header_shape.fill.solid()
    header_shape.fill.fore_color.rgb = template['title_color']
    header_shape.line.fill.background() # No border
    
    # 3. Header Text
    htf = header_shape.text_frame
    htf.text = "STRATEGIC INSIGHT"
    htf.paragraphs[0].alignment = PP_PARAGRAPH_ALIGNMENT.CENTER
    htf.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)
    htf.paragraphs[0].font.bold = True
    htf.paragraphs[0].font.size = Pt(14)
    htf.vertical_anchor = MSO_ANCHOR.MIDDLE

    # 4. Body Text
    tf = shape.text_frame
    tf.margin_top = Inches(0.8) # Push text below header
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    tf.word_wrap = True
    
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(12)
    p.font.color.rgb = template['text_color']
    p.alignment = PP_PARAGRAPH_ALIGNMENT.LEFT

def add_formatted_notes(notes_text_frame, slide_data):
    """
    Populates the speaker notes section of the slide.
    """
    tf = notes_text_frame
    tf.clear()
    
    def add_section(header, content):
        if content and content not in ["Suggestion not provided by AI.", "N/A"]:
            p = tf.add_paragraph()
            p.text = f"{header}: {content}"
            p.font.size = Pt(10)

    add_section("Speaker Notes", slide_data.get('notes'))
    add_section("Elaboration", slide_data.get('elaboration'))
    add_section("Tip", slide_data.get('best_practice_tip'))

def create_presentation(all_slides_data, template_name='professional', any_truncated=False, num_processed=0):
    """
    Generates a .pptx file in memory based on the provided slide data.
    
    Args:
        all_slides_data (dict): {filename: [slide_dict_1, slide_dict_2, ...]}
        template_name (str): Key from TEMPLATES dict.
        
    Returns:
        BytesIO: The binary PowerPoint file.
    """
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    
    t_config = TEMPLATES.get(template_name, TEMPLATES['professional'])

    # --- Title Slide ---
    slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank layout
    
    # Decorative Bar
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, Inches(2.5), prs.slide_width, Inches(2.5))
    bar.fill.solid()
    bar.fill.fore_color.rgb = t_config['title_color']
    bar.line.fill.background()

    # Title Text
    # We use the first key as the main title (usually the filename)
    source_name = list(all_slides_data.keys())[0] if all_slides_data else "Document Summary"
    
    title_box = slide.shapes.add_textbox(Inches(1), Inches(3), Inches(11.3), Inches(1.5))
    title_p = title_box.text_frame.paragraphs[0]
    title_p.text = f"Executive Summary:\n{source_name}"
    title_p.font.color.rgb = RGBColor(255, 255, 255)
    title_p.font.size = Pt(40)
    title_p.font.bold = True
    title_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # Subtitle
    sub_box = slide.shapes.add_textbox(Inches(1), Inches(5.2), Inches(11.3), Inches(0.5))
    sub_p = sub_box.text_frame.paragraphs[0]
    sub_p.text = "Generated by Trustif.ai Analyst Engine"
    sub_p.font.color.rgb = t_config['text_color']
    sub_p.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER

    # --- Content Slides ---
    for source_key, slides in all_slides_data.items():
        if not isinstance(slides, list): continue
        
        for i, s_data in enumerate(slides):
            slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
            
            # 1. Slide Title Bar
            title_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.5), Inches(0.3), Inches(12.33), Inches(0.8))
            title_bar.fill.background() # Transparent
            title_bar.line.color.rgb = t_config['accent_color']
            title_bar.line.width = Pt(2)
            
            # Title Text
            tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12), Inches(0.6))
            tp = tb.text_frame.paragraphs[0]
            tp.text = s_data.get('title', f'Slide {i+1}')
            tp.font.size = Pt(28)
            tp.font.color.rgb = t_config['title_color']
            tp.font.bold = True

            # 2. Key Message (The "So What")
            if s_data.get('key_message'):
                km_box = slide.shapes.add_textbox(Inches(0.6), Inches(1.2), Inches(8), Inches(0.6))
                km = km_box.text_frame.paragraphs[0]
                km.text = s_data.get('key_message', '')
                km.font.size = Pt(14)
                km.font.color.rgb = t_config['accent_color']
                km.font.italic = True

            # 3. Main Content (Bullets)
            body_box = slide.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(8), Inches(4.5))
            tf = body_box.text_frame
            tf.word_wrap = True
            for bullet in s_data.get('bullets', []):
                p = tf.add_paragraph()
                p.text = bullet
                p.font.size = Pt(18)
                p.font.color.rgb = t_config['text_color']
                p.space_after = Pt(12)
                p.level = 0

            # 4. Right Sidebar (Analyst Insight)
            # Checks for 'strategic_takeaway' (new agent) or fallback to 'visual' (legacy)
            insight_text = s_data.get('strategic_takeaway', s_data.get('visual', ''))
            add_insight_panel(slide, insight_text, t_config, Inches(9), Inches(2.0), Inches(4), Inches(4.5))

            # 5. Footer
            footer = slide.shapes.add_textbox(Inches(0.5), Inches(7), Inches(12.33), Inches(0.3))
            fp = footer.text_frame.paragraphs[0]
            fp.text = f"{source_key} | Page {i+1}"
            fp.font.size = Pt(9)
            fp.font.color.rgb = RGBColor(150, 150, 150)
            
            # Notes
            if slide.has_notes_slide:
                add_formatted_notes(slide.notes_slide.notes_text_frame, s_data)

    output = BytesIO()
    prs.save(output)
    output.seek(0)
    return output