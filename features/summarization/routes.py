# features/summarization/routes.py
from flask import (
    Blueprint, render_template, request, flash, current_app, url_for, g, redirect, send_file, jsonify, after_this_request, Response, stream_with_context
)
import os
import io
import uuid
import re
import logging
import time
import threading
import json
from werkzeug.utils import secure_filename
import google.generativeai as genai
from docx import Document
from pptx import Presentation
import pandas as pd
import fitz

bp = Blueprint('summarization', __name__)

# --- Load PPT Logic ---
from .ppt_builder_logic import prompts as ppt_prompts
from .ppt_builder_logic import translate as ppt_translate
from .ppt_builder_logic.file_processor import (
    allowed_file as ppt_allowed_file_util,
    extract_text_from_blob as ppt_extract_text_from_blob
)
from .ppt_builder_logic.url_processor import is_safe_url as ppt_is_safe_url, fetch_and_extract_url_content as ppt_fetch_and_extract_url_content
from .ppt_builder_logic.presentation_generator import create_presentation as ppt_create_presentation
from .ppt_builder_logic.core_processor import parse_llm_output as ppt_parse_llm_output
from .ppt_builder_logic.core_processor import generate_slides_from_text as ppt_generate_slides_from_text
PPT_BUILDER_MODULES_LOADED = True

TEXT_SUMMARY_MAX_FILE_SIZE_MB = 10 
TEXT_SUMMARY_MAX_FILE_SIZE_BYTES = TEXT_SUMMARY_MAX_FILE_SIZE_MB * 1024 * 1024
DEFAULT_PPT_MAX_CONCURRENT_REQUESTS = 3
ppt_processing_semaphore = threading.Semaphore(DEFAULT_PPT_MAX_CONCURRENT_REQUESTS)

# --- File Reading Utils ---
def read_text_from_docx(file_stream):
    try:
        file_stream.seek(0); doc = Document(file_stream)
        return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])
    except: return ""

def read_text_from_pptx(file_stream):
    try:
        file_stream.seek(0); ppt = Presentation(file_stream); runs = []
        for s in ppt.slides:
            for sh in s.shapes:
                if sh.has_text_frame: runs.extend([r.text for p in sh.text_frame.paragraphs for r in p.runs])
        return "\n".join(runs)
    except: return ""

def read_text_from_excel(file_stream):
    try:
        file_stream.seek(0); df = pd.read_excel(file_stream); return df.to_string()
    except: return ""

def read_text_from_pdf(file_stream):
    try:
        file_stream.seek(0); doc = fitz.open(stream=file_stream.read(), filetype="pdf"); return "\n".join([p.get_text() for p in doc])
    except: return ""

# --- Call LLM for PPT (Standard) ---
def call_llm_for_ppt_builder(prompt_text, max_output_tokens=8192):
    model = genai.GenerativeModel(current_app.config.get('GEMINI_MODEL_NAME'))
    response = model.generate_content(prompt_text)
    if response and response.text: return response.text
    return ""

# --- Streaming Generator for Text Summary ---
def stream_gemini_analysis(text_content, model_name, filename=""):
    try:
        model = genai.GenerativeModel(model_name)
        
        # 1. Classify
        excerpt = text_content[:5000]
        classification = "General Business Document"
        try:
            class_prompt = ppt_prompts.get_classification_prompt(excerpt)
            class_resp = model.generate_content(class_prompt)
            if class_resp and class_resp.text:
                classification = ppt_prompts.parse_classification_response(class_resp.text, filename)
        except: pass
        
        yield json.dumps({"type": "meta", "classification": classification}) + "\n"

        # 2. Summarize (With Thinking)
        prompt = ppt_prompts.build_expert_text_summary_prompt(text_content, classification)
        response_stream = model.generate_content(prompt, stream=True)
        
        for chunk in response_stream:
            if chunk.text:
                yield json.dumps({"type": "chunk", "content": chunk.text}) + "\n"
                
    except Exception as e:
        yield json.dumps({"type": "error", "content": str(e)}) + "\n"


# --- Streaming Generator for PPT Creation (New Feature) ---
def stream_ppt_creation_status(text, filename, tpl, aud, tone, lang, req_id):
    """
    Executes PPT logic step-by-step and yields status messages to the UI.
    """
    try:
        # Step 1: LLM Processing
        yield json.dumps({"type": "status", "message": "Analyzing document structure..."}) + "\n"
        
        # We manually call generate_slides logic here so we can inspect the result
        # Note: We use the synchronous generate function, but since we are in a generator,
        # the initial 'yield' has already flushed headers to the client.
        
        yield json.dumps({"type": "status", "message": "Generating slide content (this may take a moment)..."}) + "\n"
        
        slides = ppt_generate_slides_from_text(
            text, filename, call_llm_for_ppt_builder, ppt_translate.translate_slide_data,
            lang, aud, tone, tpl, False
        )

        if isinstance(slides, str) and slides.startswith("ERROR"):
             raise Exception(slides)

        # Step 2: Inspection (The Log Message You Wanted)
        count = len(slides) if isinstance(slides, list) else 0
        yield json.dumps({"type": "status", "message": f"Found {count} potential slide blocks."}) + "\n"
        time.sleep(0.5) # Slight pause so user sees the message

        # Step 3: Construction
        yield json.dumps({"type": "status", "message": f"Creating presentation for 1 source..."}) + "\n"
        pptx_buffer = ppt_create_presentation({filename: slides}, tpl)
        
        # Step 4: Finalizing
        yield json.dumps({"type": "status", "message": "Adding slide numbers and footer..."}) + "\n"
        time.sleep(0.5)

        # Step 5: Upload
        safe_name = re.sub(r'[^\w\-]+', '_', os.path.splitext(filename)[0])[:50]
        dl_name = f"{safe_name}_Deck_{req_id}.pptx"
        gcs_path = f"{req_id}/output/{dl_name}"
        
        yield json.dumps({"type": "status", "message": "Uploading to secure storage..."}) + "\n"
        
        if current_app.config.get('GCS_AVAILABLE'):
             blob = current_app.gcs_bucket.blob(gcs_path)
             pptx_buffer.seek(0)
             blob.upload_from_file(pptx_buffer, content_type='application/vnd.openxmlformats-officedocument.presentationml.presentation')

        dl_url = url_for('summarization.download_generated_ppt', file_id=req_id, filename=dl_name)
        
        # Step 6: Success HTML
        success_html = f"""
        <div style="background: #f0fdf4; border: 1px solid #bbf7d0; padding: 2rem; border-radius: 8px; text-align: center; color: #166534; animation: fadeIn 0.5s;">
            <div style="background: #dcfce7; width: 60px; height: 60px; border-radius: 50%; display: flex; align-items: center; justify-content: center; margin: 0 auto 1rem auto;">
                <i class="fas fa-check" style="font-size: 1.5rem; color: #16a34a;"></i>
            </div>
            <h3 style="margin: 0 0 0.5rem 0; color: #14532d; font-size: 1.25rem;">Presentation Ready!</h3>
            <p style="margin-bottom: 1.5rem; color: #166534;">Created {count} slides from source.</p>
            <a href="{dl_url}" class="submit-button" style="background: #16a34a; box-shadow: 0 4px 6px -1px rgba(22, 163, 74, 0.4); border: none; color: white; padding: 0.75rem 1.5rem; border-radius: 8px; text-decoration: none; display: inline-flex; align-items: center; gap: 0.5rem;">
                <i class="fas fa-download"></i> Download .pptx
            </a>
        </div>
        """
        
        yield json.dumps({"type": "result", "html": success_html}) + "\n"

    except Exception as e:
        logging.error(f"PPT Generation Error: {e}")
        yield json.dumps({"type": "error", "message": str(e)}) + "\n"


# --- Routes ---

@bp.route("/process/summarization/summarize", methods=["POST"])
def process_summarize():
    if not current_app.config.get('GEMINI_CONFIGURED'): return jsonify({"error": "Gemini not configured"}), 503
    
    text = ""
    fname = "Document"
    if request.form.get("use_sample") == "true":
        path = os.path.join(current_app.root_path, 'static', 'files', "IBM-2024-annual-report.pdf")
        if os.path.exists(path):
            with open(path, 'rb') as f: text = read_text_from_pdf(io.BytesIO(f.read()))
            fname = "IBM_Report.pdf"
    else:
        file = request.files.get("file")
        if file:
            fname = secure_filename(file.filename)
            stream = io.BytesIO(file.read())
            ext = os.path.splitext(fname)[1].lower()
            if ext == '.docx': text = read_text_from_docx(stream)
            elif ext == '.pdf': text = read_text_from_pdf(stream)
            elif ext == '.pptx': text = read_text_from_pptx(stream)
            elif ext == '.xlsx': text = read_text_from_excel(stream)

    if not text: return jsonify({"error": "No text extracted"}), 400

    return Response(stream_with_context(stream_gemini_analysis(text, current_app.config.get('GEMINI_MODEL_NAME'), fname)), mimetype='application/x-ndjson')

@bp.route("/process/summarization/create_ppt", methods=["POST"])
def process_create_ppt():
    g.request_id = uuid.uuid4().hex
    
    if not current_app.config.get('GEMINI_CONFIGURED'):
        return jsonify({"error": "Gemini not configured"}), 503

    text = ""
    fname = "Presentation"
    
    try:
        if request.form.get("use_sample") == "true":
            path = os.path.join(current_app.root_path, 'static', 'files', "IBM-2024-annual-report.pdf")
            if os.path.exists(path):
                with open(path, 'rb') as f: text = read_text_from_pdf(io.BytesIO(f.read()))
                fname = "IBM_Report.pdf"
        else:
            file = request.files.get("file")
            if file:
                fname = secure_filename(file.filename)
                stream = io.BytesIO(file.read())
                ext = os.path.splitext(fname)[1].lower()
                if ext == '.docx': text = read_text_from_docx(stream)
                elif ext == '.pdf': text = read_text_from_pdf(stream)
        
        if not text:
             return jsonify({"error": "No text extracted (PDF/DOCX only)"}), 400

        tpl = request.form.get('template', 'professional')
        aud = request.form.get('audience', 'Executives')
        tone = request.form.get('tone', 'Informative')
        lang = request.form.get('language', 'English')
        
        # Use the new Streaming Generator
        return Response(
            stream_with_context(stream_ppt_creation_status(text, fname, tpl, aud, tone, lang, g.request_id)),
            mimetype='application/x-ndjson'
        )

    except Exception as e:
        logging.error(f"PPT Error: {e}", exc_info=True)
        return jsonify({"error": str(e)}), 500

@bp.route("/download/ppt/<file_id>/<filename>", methods=["GET"])
def download_generated_ppt(file_id, filename):
    gcs_path = f"{file_id}/output/{filename}" 
    if not current_app.config.get('GCS_AVAILABLE'):
        flash("Download unavailable.", "error")
        return redirect(url_for('main.index', feature_key='summarization'))
    try:
        blob = current_app.gcs_bucket.blob(gcs_path)
        if not blob.exists():
            flash("File expired.", "error")
            return redirect(url_for('main.index', feature_key='summarization'))
        buffer = io.BytesIO()
        blob.download_to_file(buffer)
        buffer.seek(0)
        return send_file(buffer, as_attachment=True, download_name=filename, mimetype='application/vnd.openxmlformats-officedocument.presentationml.presentation')
    except Exception as e:
        return redirect(url_for('main.index', feature_key='summarization'))