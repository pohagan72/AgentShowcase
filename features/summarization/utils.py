# features/summarization/utils.py
import os  # <--- Added this missing import
import io
import logging
import fitz  # PyMuPDF
import pandas as pd
from docx import Document
from pptx import Presentation
from werkzeug.utils import secure_filename

def extract_text_from_stream(file_stream, file_extension):
    """
    Core logic to extract text from a binary stream based on file extension.
    
    Args:
        file_stream (BytesIO): The binary file content.
        file_extension (str): The file extension (e.g., '.pdf', '.docx').
        
    Returns:
        str: The extracted text content, or empty string on failure.
    """
    text = ""
    ext = file_extension.lower()

    try:
        # Ensure we are at the start of the stream
        file_stream.seek(0)

        if ext == '.docx':
            doc = Document(file_stream)
            # Extract text from paragraphs
            text_parts = [p.text for p in doc.paragraphs if p.text.strip()]
            # Extract text from tables (optional but often useful)
            for table in doc.tables:
                for row in table.rows:
                    for cell in row.cells:
                        if cell.text.strip():
                            text_parts.append(cell.text.strip())
            text = "\n".join(text_parts)

        elif ext == '.pdf':
            # Open PDF stream with PyMuPDF
            doc = fitz.open(stream=file_stream.read(), filetype="pdf")
            text_parts = []
            for page in doc:
                page_text = page.get_text()
                if page_text.strip():
                    text_parts.append(page_text)
            text = "\n".join(text_parts)

        elif ext == '.pptx':
            ppt = Presentation(file_stream)
            text_parts = []
            for slide in ppt.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text_frame") and shape.text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            # Join runs within a paragraph
                            para_text = "".join(r.text for r in paragraph.runs)
                            if para_text.strip():
                                text_parts.append(para_text)
            text = "\n".join(text_parts)

        elif ext == '.xlsx':
            # Read Excel into a pandas DataFrame and convert to string representation
            # sheet_name=None reads all sheets
            dfs = pd.read_excel(file_stream, sheet_name=None)
            text_parts = []
            for sheet_name, df in dfs.items():
                text_parts.append(f"--- Sheet: {sheet_name} ---")
                text_parts.append(df.to_string(index=False))
            text = "\n\n".join(text_parts)

        else:
            logging.warning(f"Unsupported file extension for text extraction: {ext}")

    except Exception as e:
        logging.error(f"Error extracting text from {ext} stream: {e}", exc_info=True)
        return ""

    return text

def read_text_from_file(file_storage):
    """
    Helper to process a Flask FileStorage object (from request.files).
    
    Args:
        file_storage (FileStorage): The file object from Flask.
        
    Returns:
        tuple: (extracted_text, secure_filename)
    """
    if not file_storage or not file_storage.filename:
        return "", ""

    filename = secure_filename(file_storage.filename)
    _, ext = os.path.splitext(filename)
    
    # Read file into memory stream
    stream = io.BytesIO(file_storage.read())
    
    text = extract_text_from_stream(stream, ext)
    
    return text, filename